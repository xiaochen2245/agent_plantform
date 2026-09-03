"""
演示汇报方案与答辩 PPT (PPTX / PPT) 解析器
基于 python-pptx 解析幻灯片标题层级、正文文本框、嵌入式表格
以及演讲者备注 (SPEAKER_NOTE，捕获核心设计意图与隐藏技术细节)。
"""

import io
import re
import uuid
import xml.etree.ElementTree as ET
import zipfile
from typing import Any, Dict, List, Optional, Tuple

try:
    import pptx
    from pptx.presentation import Presentation
    from pptx.shapes.table import Table as PPTXTable
    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False

from app.schemas.ast import (
    ASTBlockType,
    ASTNode,
    DocumentSourceType,
    TableCell,
    TableData,
    UnifiedDocumentAST,
)
from .base import BaseParser, EmptyDocumentError, MalformedDocumentError


class PPTXParser(BaseParser):
    """
    PPTX 汇报与答辩演示文稿解析适配器。
    提取幻灯片大纲结构、图表内容与隐式演讲者备注，保证 RAG 检索不遗漏关键补充口径。
    """

    @property
    def supported_extensions(self) -> List[str]:
        return [".pptx", ".ppt"]

    @property
    def source_type(self) -> DocumentSourceType:
        return DocumentSourceType.PPTX

    async def parse(
        self,
        content: bytes,
        file_name: str,
        tenant_id: str = "default",
        document_id: Optional[str] = None,
        **kwargs: Any,
    ) -> UnifiedDocumentAST:
        """异步解析 PPTX"""
        doc_id = document_id or f"doc_{uuid.uuid4().hex[:12]}"

        if not content:
            raise EmptyDocumentError(f"PPTX 文件 '{file_name}' 内容为空", file_name=file_name)

        if HAS_PYTHON_PPTX:
            try:
                prs = pptx.Presentation(io.BytesIO(content))
                return self._parse_with_python_pptx(prs, file_name, tenant_id, doc_id)
            except Exception as e:
                return self._parse_pure_xml_fallback(content, file_name, tenant_id, doc_id)
        else:
            return self._parse_pure_xml_fallback(content, file_name, tenant_id, doc_id)

    def _parse_with_python_pptx(
        self,
        prs: Any,
        file_name: str,
        tenant_id: str,
        doc_id: str
    ) -> UnifiedDocumentAST:
        """基于 python-pptx 逐页提取"""
        nodes: List[ASTNode] = []
        total_slides = len(prs.slides)

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_label = f"Slide {slide_idx}"
            
            # 1. 抽取幻灯片标题
            slide_title = ""
            if slide.shapes.title and slide.shapes.title.text:
                slide_title = self.clean_text(slide.shapes.title.text)

            if not slide_title:
                slide_title = f"第 {slide_idx} 页演示方案"

            title_node = ASTNode(
                block_id=self.generate_block_id("pptx_title"),
                block_type=ASTBlockType.HEADING,
                level=1,
                section_path=[slide_title],
                text_content=slide_title,
                page_or_sheet=slide_label,
                extra_metadata={"slide_index": slide_idx}
            )
            nodes.append(title_node)

            # 2. 空间几何排序抽取形状与文本框
            # 过滤掉已经是 title 的 shape，按 top (Y), left (X) 排序
            content_shapes = []
            for s in slide.shapes:
                if s == slide.shapes.title:
                    continue
                content_shapes.append(s)

            content_shapes.sort(key=lambda s: (getattr(s, "top", 0) or 0, getattr(s, "left", 0) or 0))

            for shape in content_shapes:
                # 检查表格
                if shape.has_table:
                    table_node = self._parse_pptx_table(shape.table, slide_title, slide_label)
                    if table_node:
                        nodes.append(table_node)
                # 检查普通文本
                elif shape.has_text_frame:
                    t_text = self.clean_text(shape.text_frame.text)
                    if t_text:
                        inferred_lvl = self.infer_heading_level(t_text)
                        b_type = ASTBlockType.HEADING if inferred_lvl else ASTBlockType.PARAGRAPH
                        nodes.append(
                            ASTNode(
                                block_id=self.generate_block_id("pptx_txt"),
                                block_type=b_type,
                                level=inferred_lvl or 2,
                                section_path=[slide_title, t_text[:30]],
                                text_content=t_text,
                                page_or_sheet=slide_label
                            )
                        )

            # 3. 提取演讲者备注 (Speaker Notes)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = self.clean_text(slide.notes_slide.notes_text_frame.text)
                # 过滤默认占位符或空白备注
                if notes_text and not notes_text.startswith("Click to edit Master"):
                    nodes.append(
                        ASTNode(
                            block_id=self.generate_block_id("pptx_note"),
                            block_type=ASTBlockType.SPEAKER_NOTE,
                            section_path=[slide_title, "演讲者备注"],
                            text_content=notes_text,
                            page_or_sheet=slide_label,
                            extra_metadata={"is_speaker_note": True}
                        )
                    )

        if not nodes:
            nodes.append(
                ASTNode(
                    block_id=self.generate_block_id("pptx_empty"),
                    block_type=ASTBlockType.PARAGRAPH,
                    text_content="[PPT 演示文稿未包含有效文本或表格]",
                    page_or_sheet="Slide 1"
                )
            )

        return UnifiedDocumentAST(
            document_id=doc_id,
            tenant_id=tenant_id,
            file_name=file_name,
            source_type=self.source_type,
            total_pages_or_sheets=max(total_slides, 1),
            nodes=nodes,
            metadata={
                "parser": "PPTXParser",
                "total_slides": total_slides,
                "extracted_nodes_count": len(nodes),
            }
        )

    def _parse_pptx_table(
        self,
        table: Any,
        slide_title: str,
        slide_label: str
    ) -> Optional[ASTNode]:
        """解析 PPT 嵌入表格"""
        if not table.rows:
            return None

        n_rows = len(table.rows)
        n_cols = len(table.columns)
        if n_cols == 0:
            return None

        grid: List[List[str]] = []
        cells: List[TableCell] = []

        for r_idx, row in enumerate(table.rows):
            row_vals: List[str] = []
            for c_idx, cell in enumerate(row.cells):
                ct = self.clean_text(cell.text)
                row_vals.append(ct)
                cells.append(
                    TableCell(
                        row=r_idx,
                        col=c_idx,
                        row_span=1,
                        col_span=1,
                        text=ct,
                        is_header=(r_idx == 0)
                    )
                )
            grid.append(row_vals)

        headers = [grid[0]] if grid else []
        body = grid[1:] if len(grid) > 1 else []
        md = self.build_table_markdown(headers=headers, rows=body)

        return ASTNode(
            block_id=self.generate_block_id("pptx_tbl"),
            block_type=ASTBlockType.TABLE,
            level=2,
            section_path=[slide_title, "演示表格"],
            text_content=md,
            table_data=TableData(
                headers=headers,
                rows=body,
                cells=cells,
                markdown=md,
                summary=f"PPT 嵌入表格，共 {len(body)} 行。"
            ),
            page_or_sheet=slide_label
        )

    def _parse_pure_xml_fallback(
        self,
        content: bytes,
        file_name: str,
        tenant_id: str,
        doc_id: str
    ) -> UnifiedDocumentAST:
        """纯 Python 遍历 ppt/slides/ 与 ppt/notesSlides/ 降级引擎"""
        nodes: List[ASTNode] = []
        try:
            with zipfile.ZipFile(io.BytesIO(content), "r") as z:
                slide_files = sorted(
                    [n for n in z.namelist() if re.match(r"^ppt/slides/slide\d+\.xml$", n)],
                    key=lambda x: int(re.search(r"\d+", x).group(0))
                )
                
                a_ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                        "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}

                for idx, sf in enumerate(slide_files, start=1):
                    slide_label = f"Slide {idx}"
                    slide_xml = z.read(sf)
                    try:
                        root = ET.fromstring(slide_xml)
                    except Exception:
                        continue

                    # 提取标题
                    slide_title = ""
                    for sp in root.findall(".//p:sp", a_ns):
                        ph = sp.find(".//p:ph", a_ns)
                        if ph is not None and ph.get("type") in ("title", "ctrTitle"):
                            t_elems = sp.findall(".//a:t", a_ns)
                            slide_title = self.clean_text("".join(t.text or "" for t in t_elems))
                            break

                    if not slide_title:
                        # 尝试首个文本
                        first_t = root.find(".//a:t", a_ns)
                        if first_t is not None and first_t.text:
                            slide_title = self.clean_text(first_t.text)
                        else:
                            slide_title = f"第 {idx} 页演示方案"

                    nodes.append(
                        ASTNode(
                            block_id=self.generate_block_id("pptx_title_fb"),
                            block_type=ASTBlockType.HEADING,
                            level=1,
                            section_path=[slide_title],
                            text_content=slide_title,
                            page_or_sheet=slide_label,
                            extra_metadata={"slide_index": idx}
                        )
                    )

                    # 提取表格 <a:tbl>
                    for tbl in root.findall(".//a:tbl", a_ns):
                        tr_elems = tbl.findall(".//a:tr", a_ns)
                        grid_rows: List[List[str]] = []
                        cells: List[TableCell] = []
                        for r_i, tr in enumerate(tr_elems):
                            tc_elems = tr.findall(".//a:tc", a_ns)
                            r_vals: List[str] = []
                            for c_i, tc in enumerate(tc_elems):
                                c_t = self.clean_text("".join(t.text or "" for t in tc.findall(".//a:t", a_ns)))
                                r_vals.append(c_t)
                                cells.append(
                                    TableCell(row=r_i, col=c_i, row_span=1, col_span=1, text=c_t, is_header=(r_i == 0))
                                )
                            grid_rows.append(r_vals)
                        if grid_rows:
                            headers = [grid_rows[0]]
                            body = grid_rows[1:] if len(grid_rows) > 1 else []
                            md = self.build_table_markdown(headers=headers, rows=body)
                            nodes.append(
                                ASTNode(
                                    block_id=self.generate_block_id("pptx_tbl_fb"),
                                    block_type=ASTBlockType.TABLE,
                                    level=2,
                                    section_path=[slide_title, "演示表格"],
                                    text_content=md,
                                    table_data=TableData(
                                        headers=headers,
                                        rows=body,
                                        cells=cells,
                                        markdown=md,
                                        summary=f"PPT 嵌入表格，共 {len(body)} 行。"
                                    ),
                                    page_or_sheet=slide_label
                                )
                            )

                    # 提取普通形状与组合形状文本（带坐标过滤）
                    for sp in root.findall(".//p:sp", a_ns):
                        # 检查 off-canvas 坐标
                        off = sp.find(".//a:off", a_ns)
                        if off is not None:
                            try:
                                x = int(off.get("x", "0"))
                                y = int(off.get("y", "0"))
                                if x < 0 or y < 0:
                                    continue  # 过滤离屏草稿对象
                            except ValueError:
                                pass

                        # 过滤已作为 title 的形状
                        ph = sp.find(".//p:ph", a_ns)
                        if ph is not None and ph.get("type") in ("title", "ctrTitle"):
                            continue

                        texts = [t.text.strip() for t in sp.findall(".//a:t", a_ns) if t.text and t.text.strip()]
                        t_text = self.clean_text("\n".join(texts))
                        if t_text and t_text != slide_title:
                            nodes.append(
                                ASTNode(
                                    block_id=self.generate_block_id("pptx_txt_fb"),
                                    block_type=ASTBlockType.PARAGRAPH,
                                    section_path=[slide_title],
                                    text_content=t_text,
                                    page_or_sheet=slide_label
                                )
                            )

                    # 提取演讲者备注: ppt/notesSlides/notesSlide{idx}.xml
                    notes_path = f"ppt/notesSlides/notesSlide{idx}.xml"
                    if notes_path in z.namelist():
                        notes_xml = z.read(notes_path)
                        n_root = ET.fromstring(notes_xml)
                        n_texts = [t.text.strip() for t in n_root.findall(".//a:t", a_ns) if t.text and t.text.strip()]
                        n_full = self.clean_text("\n".join(n_texts))
                        if n_full and not n_full.startswith("Click to edit Master"):
                            nodes.append(
                                ASTNode(
                                    block_id=self.generate_block_id("pptx_note_fb"),
                                    block_type=ASTBlockType.SPEAKER_NOTE,
                                    section_path=[slide_title, "演讲者备注"],
                                    text_content=n_full,
                                    page_or_sheet=slide_label,
                                    extra_metadata={"is_speaker_note": True}
                                )
                            )
        except zipfile.BadZipFile as e:
            raise MalformedDocumentError(f"解压 PPTX 归档失败: {e}", file_name=file_name) from e
        except Exception as e:
            if isinstance(e, (EmptyDocumentError, MalformedDocumentError)):
                raise
            raise MalformedDocumentError(f"解析 PPTX 失败: {e}", file_name=file_name) from e

        return UnifiedDocumentAST(
            document_id=doc_id,
            tenant_id=tenant_id,
            file_name=file_name,
            source_type=self.source_type,
            total_pages_or_sheets=max(len(slide_files) if 'slide_files' in locals() else 1, 1),
            nodes=nodes,
            metadata={"parser": "PPTXParser", "is_pure_xml_fallback": True, "extracted_nodes_count": len(nodes)}
        )
